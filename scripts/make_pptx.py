"""
Generate the final hackathon slide deck as PPTX.

Usage: python scripts/make_pptx.py

Output: outputs/slides/Chirality_Atlas_Star_Ascidian_Edition.pptx
"""
import sys
import os
sys.path.insert(0, os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))), "src"))

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OUT_DIR = os.path.join(REPO, "outputs", "slides")
os.makedirs(OUT_DIR, exist_ok=True)

PPTX_PATH = os.path.join(OUT_DIR, "Chirality_Atlas_Star_Ascidian_Edition.pptx")

PANELS = os.path.join(REPO, "outputs", "panels")
REF = os.path.join(REPO, "assets", "reference", "star_ascidian_reference.jpg")

INK = RGBColor(0x1F, 0x24, 0x21)
ACCENT = RGBColor(0xC1, 0x5A, 0x3A)
GREEN = RGBColor(0x31, 0x5C, 0x4C)
LIGHT = RGBColor(0xF7, 0xF3, 0xEA)
BORDER = RGBColor(0xDD, 0xD5, 0xC8)
PANEL_BG = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x0D, 0x11, 0x17)
BLUE_GRAY = RGBColor(0x7A, 0xA8, 0xC8)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def _set_bg(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb


def _txt(slide, text, left, top, width, height,
         size=18, bold=False, color=None, align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color if color else INK


def _rect(slide, left, top, width, height, fill_rgb, line_rgb=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
    else:
        shape.line.fill.background()
    return shape


def _img(slide, path, left, top, width, height=None):
    if not os.path.exists(path):
        return False
    if height:
        slide.shapes.add_picture(path, left, top, width, height)
    else:
        slide.shapes.add_picture(path, left, top, width)
    return True


def _slide_footer(slide, n, total=5):
    _txt(slide, f"{n} / {total}",
         Inches(12.8), Inches(7.15), Inches(0.5), Inches(0.3),
         size=9, color=BORDER, align=PP_ALIGN.RIGHT)
    _txt(slide, "Chirality Atlas: Star Ascidian Edition",
         Inches(0.3), Inches(7.15), Inches(8.0), Inches(0.3),
         size=9, color=BORDER, align=PP_ALIGN.LEFT, italic=True)


def _title_bar(slide, title, subtitle=None):
    _rect(slide, Inches(0), Inches(0), Inches(13.33), Inches(1.0), INK)
    _txt(slide, title,
         Inches(0.35), Inches(0.08), Inches(12.6), Inches(0.65),
         size=26, bold=True, color=LIGHT, align=PP_ALIGN.LEFT)
    if subtitle:
        _txt(slide, subtitle,
             Inches(0.35), Inches(0.7), Inches(12.6), Inches(0.28),
             size=11, color=ACCENT, align=PP_ALIGN.LEFT, italic=True)


def add_slide_notes(slide, notes_text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = notes_text


def make_slide1(prs):
    """Hook: real organism + simulation match side by side."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, LIGHT)
    _title_bar(slide,
               "Can local rules generate a living star pattern?",
               "Observe: star ascidian colony  |  Hypothesize: Turing field + angular repulsion  |  Result: yes for the core geometry")

    content_top = Inches(1.10)
    content_h = Inches(5.30)   # shorter to leave room for message strip

    # Left half: real reference image (prominent)
    if os.path.exists(REF):
        _img(slide, REF, Inches(0.2), content_top, Inches(6.3), content_h)
        _txt(slide, "Botryllus schlosseri (real colony)",
             Inches(0.2), Inches(6.47), Inches(6.3), Inches(0.22),
             size=10, color=ACCENT, align=PP_ALIGN.CENTER, italic=True)

    # Right half: colony_scale_reference_match (simulation match)
    colony_path = os.path.join(PANELS, "colony_scale_reference_match.png")
    _img(slide, colony_path, Inches(6.7), content_top, Inches(6.4), content_h)
    _txt(slide, "Simulation: colony-scale view (3x3 periodic tiling)",
         Inches(6.7), Inches(6.47), Inches(6.4), Inches(0.22),
         size=10, color=BLUE_GRAY, align=PP_ALIGN.CENTER, italic=True)

    # Key message strip -- visible at bottom of content area
    _rect(slide, Inches(0.2), Inches(6.75), Inches(12.9), Inches(0.28), LIGHT, BORDER)
    _txt(slide,
         "Two-level pattern: star spacing + radial zooid geometry.  "
         "We model the geometry, not the full organism.",
         Inches(0.35), Inches(6.78), Inches(12.6), Inches(0.24),
         size=11, bold=False, color=INK, align=PP_ALIGN.CENTER)

    _slide_footer(slide, 1)

    add_slide_notes(slide,
        "Open with the real organism. Point to the star structure: shared central atrium, "
        "radial arms, regular spacing between stars. "
        "The right panel is the simulation -- many star systems tiling the domain at comparable spacing. "
        "Say: We asked whether two local rules can reproduce this geometry from scratch. "
        "The answer is yes. Let me explain how. "
        "Timing: 45-60 seconds. Transition: 'So let me walk through the mechanism.'")


def make_slide2(prs):
    """Mechanism: center_selection_schematic (Layer 1) + single_star_mechanism (Layer 2)."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, LIGHT)
    _title_bar(slide,
               "Model: centers first, stars second",
               "Layer 1: Gierer-Meinhardt Turing field  |  Layer 2: Active zooid agents with angular repulsion")

    content_top = Inches(1.10)
    content_h = Inches(5.08)   # shorter to leave room for mechanism arrows

    # Left panel: Layer 1 - center selection schematic
    center_path = os.path.join(PANELS, "center_selection_schematic.png")
    _img(slide, center_path, Inches(0.2), content_top, Inches(6.3), content_h)
    _txt(slide, "Layer 1: Turing field places star centers",
         Inches(0.2), Inches(6.25), Inches(6.3), Inches(0.25),
         size=10, color=INK, align=PP_ALIGN.CENTER)

    # Right panel: Layer 2 - single star mechanism (clean vs chiral)
    single_path = os.path.join(PANELS, "single_star_mechanism.png")
    _img(slide, single_path, Inches(6.7), content_top, Inches(6.4), content_h)
    _txt(slide, "Layer 2: angular repulsion forms discrete arms. Chirality twists.",
         Inches(6.7), Inches(6.25), Inches(6.4), Inches(0.25),
         size=10, color=INK, align=PP_ALIGN.CENTER)

    # Concise mechanism arrow strip
    mech_items = [
        ("Turing field", "centers", GREEN),
        ("radial spring + angular repulsion", "discrete arms", INK),
        ("chirality omega", "measurable swirl", ACCENT),
    ]
    _rect(slide, Inches(0.2), Inches(6.58), Inches(12.9), Inches(0.45),
          RGBColor(0xF0, 0xF0, 0xE8), BORDER)
    for i, (cause, effect, col) in enumerate(mech_items):
        x = Inches(0.4 + i * 4.3)
        _txt(slide, cause + "  ->  " + effect,
             x, Inches(6.65), Inches(4.1), Inches(0.32),
             size=11, bold=(i == 0), color=col, align=PP_ALIGN.LEFT)

    _slide_footer(slide, 2)

    add_slide_notes(slide,
        "Walk through both layers. "
        "Left: the GM field produces quasi-periodic Turing spots. Those spots are the star centers. "
        "No explicit repulsion between centers needed -- spacing comes from the diffusion ratio. "
        "Right: within each star, zooid agents are confined at r_target by a radial spring. "
        "The key ingredient is angular repulsion between arm groups. "
        "Without it: agents form a ring. With it: they separate into discrete arms. "
        "Clean star (left inset) vs chiral star (right inset) -- omega adds a turning bias. "
        "Timing: 50-60 seconds. Transition: 'Here is what this looks like in motion.'")


def make_slide3(prs):
    """Simulation dynamics: formation sequence + clean vs chiral."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, LIGHT)
    _title_bar(slide,
               "Simulation: from noise to star systems",
               "400 steps. radial_order >= 0.8. swirl_score rises from 0.01 to 0.3 at omega = 2.5.")

    # Main: formation sequence or slide3 panel
    seq_path = os.path.join(PANELS, "formation_sequence_strong.png")
    s3_path = os.path.join(PANELS, "slide3_simulation_sequence.png")
    main_path = s3_path if os.path.exists(s3_path) else seq_path
    _img(slide, main_path, Inches(0.2), Inches(1.12), Inches(12.9), Inches(5.35))

    # Key results row
    metrics = [
        ("radial_order = 1.0", "agents at r_target", GREEN),
        ("swirl = 0.01 -> 0.3", "omega 0 to 2.5", ACCENT),
        ("fragmentation < 0.1", "arms stable", GREEN),
    ]
    for i, (val, label, col) in enumerate(metrics):
        bx = Inches(0.4 + i * 4.3)
        _rect(slide, bx, Inches(6.55), Inches(4.1), Inches(0.55), PANEL_BG, BORDER)
        _txt(slide, val, bx + Inches(0.1), Inches(6.57), Inches(3.9), Inches(0.25),
             size=12, bold=True, color=col, align=PP_ALIGN.LEFT)
        _txt(slide, label, bx + Inches(0.1), Inches(6.82), Inches(3.9), Inches(0.2),
             size=9, color=INK, align=PP_ALIGN.LEFT)

    _txt(slide,
         "Top: omega = 0, radial. Bottom: omega = 2.5, chiral.  "
         "Random init -> radial arms -> chirality adds swirl.  "
         "Play live: outputs/movies/star_formation_clean.gif",
         Inches(0.4), Inches(7.13), Inches(12.5), Inches(0.25),
         size=9, color=RGBColor(0x44, 0x44, 0x44), italic=True, align=PP_ALIGN.CENTER)

    _slide_footer(slide, 3)

    add_slide_notes(slide,
        "These panels show the time evolution. "
        "At t=0 agents start in arm groups with random offsets. "
        "The radial spring drives them outward; angular repulsion pushes arms apart. "
        "Top row: omega=0. Arms settle into radial lobes. "
        "Bottom row: omega=2.5. Same initialization but arms slowly rotate. "
        "Swirl score rises from 0.01 to 0.3. The arm structure is preserved. "
        "If the app is running: Movie Gallery tab, play star_formation_clean.gif. "
        "Timing: 50-60 seconds. Transition: 'Let me show you the parameter space.'")


def make_slide4(prs):
    """Phase diagram with regime labels."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, LIGHT)
    _title_bar(slide,
               "Creative exploration: phases of living geometry",
               "4 sweeps. Named regimes. star-likeness and swirl are independent axes.")

    panel = os.path.join(PANELS, "phase_diagram_with_regimes.png")
    _img(slide, panel, Inches(0.2), Inches(1.10), Inches(12.9), Inches(5.3))

    regimes = [
        ("Uniform mat", "low Dh/Da, no Turing spots", INK),
        ("Clean stars", "moderate k_r, low omega  (best)", GREEN),
        ("Twisted stars", "moderate k_r, high omega", ACCENT),
        ("Spots, no arms", "k_r too low", INK),
        ("Merged", "overcrowded", ACCENT),
        ("Fragmented", "high Dr, noise breaks arms", ACCENT),
    ]
    for i, (regime, desc, col) in enumerate(regimes):
        col_i = i % 3
        row_i = i // 3
        x = Inches(0.4 + col_i * 4.3)
        y = Inches(6.48 + row_i * 0.30)
        _txt(slide, f"{regime}: {desc}",
             x, y, Inches(4.1), Inches(0.28),
             size=10, color=col, align=PP_ALIGN.LEFT,
             bold=(regime in ("Clean stars",)))

    # Punchline
    _txt(slide, "One run is a demo.  A phase diagram is the result.",
         Inches(0.4), Inches(7.12), Inches(12.5), Inches(0.25),
         size=11, bold=True, color=INK, align=PP_ALIGN.CENTER)

    _slide_footer(slide, 4)

    add_slide_notes(slide,
        "Phase diagram for radial spring strength versus chirality rate. "
        "Left heatmap: star-likeness. Bright at moderate k_radial, low omega -- the clean star regime. "
        "Right heatmap: swirl score rises with omega, roughly independent of k_radial up to a threshold. "
        "Key point: the two metrics are not correlated. "
        "Good arms exist without chirality. Chirality up to omega=2 does not destroy arms. "
        "Inset thumbnails show representative presets from actual simulations. "
        "Timing: 50-60 seconds. Transition: 'I want to be honest about what this model does and does not do.'")


def make_slide5(prs):
    """Insight, limits, and LLM use in three clean columns."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    _set_bg(slide, LIGHT)
    _title_bar(slide,
               "Insight, limits, and LLM use",
               "Local rules can make living geometry. Toy geometric model, not a full developmental model.")

    col_w = Inches(4.1)
    col_h = Inches(5.1)
    tops = Inches(1.12)
    gaps = [Inches(0.2), Inches(4.5), Inches(8.8)]

    # Column 1: What matched
    _rect(slide, gaps[0], tops, col_w, col_h, RGBColor(0xE8, 0xF5, 0xEC), GREEN)
    _txt(slide, "What matched", gaps[0] + Inches(0.15), tops + Inches(0.1), col_w - Inches(0.3), Inches(0.3),
         size=13, bold=True, color=GREEN)
    matched = [
        "Repeated center spacing",
        "Radial arm confinement",
        "Discrete arm structure",
        "Measurable chirality (swirl)",
        "Star-like colony tiling",
    ]
    y = tops + Inches(0.5)
    for item in matched:
        _txt(slide, "+ " + item, gaps[0] + Inches(0.15), y, col_w - Inches(0.3), Inches(0.32),
             size=11, color=INK)
        y += Inches(0.38)

    # Column 2: What did not
    _rect(slide, gaps[1], tops, col_w, col_h, RGBColor(0xFB, 0xF0, 0xEC), ACCENT)
    _txt(slide, "What did not", gaps[1] + Inches(0.15), tops + Inches(0.1), col_w - Inches(0.3), Inches(0.3),
         size=13, bold=True, color=ACCENT)
    limits = [
        "n_arms = 7 is a parameter",
        "No Botryllus biochemistry",
        "No 3D geometry",
        "No developmental staging",
        "Arm count metric underestimates",
    ]
    y = tops + Inches(0.5)
    for item in limits:
        _txt(slide, "- " + item, gaps[1] + Inches(0.15), y, col_w - Inches(0.3), Inches(0.32),
             size=11, color=INK)
        y += Inches(0.38)

    # Column 3: LLM contributions
    _rect(slide, gaps[2], tops, col_w, col_h, RGBColor(0xEA, 0xF0, 0xF7), BLUE_GRAY)
    _txt(slide, "LLM contribution", gaps[2] + Inches(0.15), tops + Inches(0.1), col_w - Inches(0.3), Inches(0.3),
         size=13, bold=True, color=BLUE_GRAY)
    llm_items = [
        "Proposed two-layer split:",
        "  Turing field for centers,",
        "  agents for arms",
        "Designed anti-cheat metrics:",
        "  radial_order, swirl_score,",
        "  fragmentation",
        "Both verified, not accepted blindly",
    ]
    y = tops + Inches(0.5)
    for item in llm_items:
        _txt(slide, item, gaps[2] + Inches(0.15), y, col_w - Inches(0.3), Inches(0.32),
             size=10, color=INK)
        y += Inches(0.35)

    # Takeaway bar
    _rect(slide, Inches(0.2), Inches(6.35), Inches(12.9), Inches(0.75),
          RGBColor(0x1F, 0x24, 0x21))
    _txt(slide,
         "Local rules can make living geometry.  "
         "Turing field places centers. Angular repulsion forms arms. "
         "Chirality adds swirl. Toy geometric model -- not a full Botryllus model.",
         Inches(0.5), Inches(6.45), Inches(12.4), Inches(0.55),
         size=12, bold=False, color=LIGHT, align=PP_ALIGN.CENTER)

    _slide_footer(slide, 5)

    add_slide_notes(slide,
        "Three columns: what matched, what did not, what LLM changed. "
        "Green: center spacing, radial confinement, discrete arms, measurable chirality. "
        "Orange: arm count is a parameter not emergent, no biochemistry, 2D only, no staging. "
        "Blue: LLM proposed the two-layer decomposition -- that was not the initial design. "
        "We started with a single-layer model. LLM suggested separating Turing-based center selection "
        "from local agent arm formation. That separation is physically principled. "
        "LLM also designed the metric hierarchy -- each metric catches a specific failure mode. "
        "Both were verified, not accepted blindly. "
        "End: Local rules can make living geometry. "
        "This is a toy geometric model -- not a full Botryllus developmental model, and we say so. "
        "Timing: 60-70 seconds.")


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    print("Building slides...")
    make_slide1(prs)
    print("  Slide 1 done")
    make_slide2(prs)
    print("  Slide 2 done")
    make_slide3(prs)
    print("  Slide 3 done")
    make_slide4(prs)
    print("  Slide 4 done")
    make_slide5(prs)
    print("  Slide 5 done")

    prs.save(PPTX_PATH)
    size_kb = os.path.getsize(PPTX_PATH) // 1024
    print(f"\nSaved: {PPTX_PATH}")
    print(f"Size:  {size_kb} KB")
    print()
    print("Slides use these assets:")
    for fname, desc in [
        ("assets/reference/star_ascidian_reference.jpg", "Slide 1 left"),
        ("outputs/panels/colony_scale_reference_match.png", "Slide 1 right"),
        ("outputs/panels/center_selection_schematic.png", "Slide 2 left"),
        ("outputs/panels/single_star_mechanism.png", "Slide 2 right"),
        ("outputs/panels/slide3_simulation_sequence.png", "Slide 3 main"),
        ("outputs/panels/phase_diagram_with_regimes.png", "Slide 4 main"),
        ("outputs/panels/slide5_insight_and_limits.png", "Slide 5 (unused, replaced by 3-column layout)"),
    ]:
        path = os.path.join(REPO, fname)
        exists = "OK" if os.path.exists(path) else "MISSING"
        print(f"  [{exists}] {desc}: {fname}")


if __name__ == "__main__":
    main()
